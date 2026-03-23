"""
=============================================================================
  AI FOR BUSINESS — Week 07
  Topic : Public relations, communications, and AI
  File  : 07.py

  Approach : AI-driven Presentation Creation
    1. Create the title slide.
    2. Add a content slide.
    3. Add an image slide (branding).

  Dependencies : python-pptx  (pip install python-pptx)
=============================================================================
"""

from pptx import Presentation 
from pptx.util import Inches 
from PIL import Image 
import os 
#  Create output folder if not exists 
os.makedirs("output", exist_ok=True) 
# Create presentation 
prs = Presentation() 
# Slide 1 Creation 
slide_layout = prs.slide_layouts[0] 
slide = prs.slides.add_slide(slide_layout) 
slide.shapes.title.text = "AI in Public Relations" 
slide.placeholders[1].text = "Automating PR and Branding with Python" 
# Slide 2 Creation 
slide_layout = prs.slide_layouts[1] 
slide = prs.slides.add_slide(slide_layout) 
slide.shapes.title.text = "Applications of AI in PR" 
tf = slide.placeholders[1].text_frame 
tf.text = "Media Monitoring" 
for point in ["Chatbots", "Sentiment Analysis", "Content Generation"]: p = tf.add_paragraph() 
p.text = point 
# Slide 3 Creation 
slide_layout = prs.slide_layouts[5] 
slide = prs.slides.add_slide(slide_layout) 
slide.shapes.title.text = "Brand Logo" 
img_path = "assets/logo.JPG" 
try: 
  img = Image.open(img_path) 
  img.verify() 
  slide.shapes.add_picture( 
  img_path, 
  Inches(2), 
  Inches(2), 
  width=Inches(4) 
) 
except Exception as e: 
  print(" Image error:", e) 
# Save in OUTPUT folder 
output_file = "output/final_presentation.pptx" 
prs.save(output_file) 
print(" PPT saved successfully in OUTPUT folder!") 
print(" Location:", output_file)